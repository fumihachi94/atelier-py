# encoding: utf-8

import pptxBasics
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import PP_ALIGN

def title_page(prs, data , slide_idx, sld_layout_idx=0):
    print("create title page")
    title_left   = Inches(1.0)
    title_buttom = Inches(2.25)
    auther = ""
    date = ""

    slide_layout = prs.slide_layouts[sld_layout_idx]
    if sld_layout_idx==0 and len(prs.slides)>0:
        slide = prs.slides[sld_layout_idx]
    elif len(prs.slides)>slide_idx:
        slide = prs.slides[slide_idx]
    else:
        slide = prs.slides.add_slide(slide_layout)

    for datum in data:
        if datum.startswith('# '):
            placeholder = slide.placeholders[0]
            placeholder.text = datum.replace('# ', '')
            placeholder.width = Inches(prs.slide_width.inches-placeholder.left.inches*2)
            placeholder.height = Inches(1.6)
            title_left   = placeholder.left.inches
            title_buttom = placeholder.top.inches+placeholder.height.inches
        elif datum.startswith('## '):
            placeholder = slide.placeholders[0]
            placeholder.text = datum.replace('## ', '')
            placeholder.width = Inches(prs.slide_width.inches-placeholder.left.inches*2)
            placeholder.height = Inches(1.25)
        elif datum.startswith('subtitle: '):
            placeholder = slide.placeholders[1]
            placeholder.text = datum.replace('subtitle: ', '')
            placeholder.left = Inches(title_left)
            placeholder.top  = Inches(title_buttom)+Inches(0.2)
            slide.shapes.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        elif datum.startswith('author: '):
            auther = datum.replace('author: ', '')
        elif datum.startswith('date: '):
            date = datum.replace('date: ', '')
        else:
            continue

    # placeholder以外のテキストBoxを削除
    for shape in slide.shapes:
        if not shape.shape_type == 14 and shape.has_text_frame:
            shape.text_frame.clear()
    
    if not date=="" or not auther == "":
        signature = slide.shapes.add_textbox(Inches(title_left*1.5), prs.slide_height*0.7, Inches(1), Inches(1))
        signature.text = date + "\n"+ auther


if __name__ == "__main__":

    filename = 'test.pptx'
    prs, file_exist = pptxBasics.open_pptx(filename)

    # Set slide size
    prs.slide_height=Inches(9)
    prs.slide_width=Inches(16)
    
    f = open('sample.md', 'r', encoding='UTF-8')

    SLD_LAYOUT_NUMBER = 0
    slide_num = 0
    data = []
    for line in f:
        line = line.rstrip('\n')

        if line.startswith('# '):
            SLD_LAYOUT_NUMBER = 0
        elif line.startswith('## '):
            SLD_LAYOUT_NUMBER = 1

        if line.startswith('---'):
            title_page(prs, data, slide_num, SLD_LAYOUT_NUMBER)
            data = []
            slide_num = slide_num+1
        else:
            data.append(line)

    f.close()

    # Set title allignment
    for sld in prs.slides:
        if sld.shapes.title:
            sld.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    pptxBasics.save_pptx(filename, prs, file_exist)