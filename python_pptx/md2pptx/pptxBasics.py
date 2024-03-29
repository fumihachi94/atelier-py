import os, sys
import datetime
from pptx import Presentation
from tkinter import messagebox

def open_pptx(filename):
    '''
    search target pptx file and open. If don't exit target file, create new presentation

    [input]
    - filename : should be "*.pptx" file
    ---
    [output]
    - 1: Presentation() Object
    - 2: File exist (0: don't exist, 1: exist)
    '''
    if not filename.endswith('.pptx'):
        print('Error: Please specify ".pptx" file!')
        sys.exit()

    filepath = os.getcwd()+"/"+filename
    
    if os.path.exists(filepath):
        return Presentation(filepath), 1
    else:
        return Presentation(), 0

def yes_no_messagebox(filename):
    return messagebox.askyesno('ファイルの保存', \
    filename+'は既に存在します。上書きしますか？ (Noを選択すると別名で保存されます)')

def save_pptx(filename, pptx_obj, file_exist):
    try:
        if file_exist:
            if yes_no_messagebox(filename):
                pptx_obj.save(filename)
            else:
                dt_now = datetime.datetime.now()
                name   = os.path.splitext(os.path.basename(filename))[0]
                pptx_obj.save(name+dt_now.strftime('_%Y%m%d_%H%M%S')+'.pptx')
        else:
            pptx_obj.save(filename)
        print("Complete file saving!")
    except:
        print("Error: can't save .pptx file.")