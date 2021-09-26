import pptxBasics
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN

filename = 'slide_layout.pptx'
prs, file_exist = pptxBasics.open_pptx(filename)

# # Set slide size
# prs.slide_height=Inches(9)
# prs.slide_width=Inches(16)


'''
0: Title (presentation title slide)
1: Title and Content
2: Section Header (sometimes called Segue)
3: Two Content (side by side bullet textboxes)
4: Comparison (same but additional title for each side by side content box)
5: Title Only
6: Blank
7: Content with Caption
8: Picture with Caption
'''
for sld_layout_idx in range(9):
    slide_layout = prs.slide_layouts[sld_layout_idx]
    if sld_layout_idx < len(prs.slides):
        slide = prs.slides[sld_layout_idx]
    else:
        slide = prs.slides.add_slide(slide_layout)

    print("---"+slide_layout.name+"---")
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))
        placeholder = slide.placeholders[shape.placeholder_format.idx]
        placeholder.text = shape.name
        print('left: %.2f inches' % placeholder.left.inches)
    print("")

# Set title allignment
# for sld in prs.slides:
#     if sld.shapes.title:
#         sld.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
 
pptxBasics.save_pptx(filename, prs, 0)
