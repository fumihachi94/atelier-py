from datetime import date
import os, sys
from typing import Awaitable, Text
import pptxBasics
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt, Inches

def create_shape(slide, shape_type, position):
    '''
    Parameters
    ----------
    slide: pptx.slide.Slide
       Presentation slide obj
    shape_type: pptx.enum.base.EnumValue
       MSO_SHAPE type
    position: struct of Cm, Inchis, Pt
       [left, top, width , height]

    Returns
    -------
    shape: pptx.shapes.autoshape.Shape
    '''
    shapes = slide.shapes
    shape = shapes.add_shape(
        shape_type, position[0], position[1], position[2], position[3]
    )

    return shape

class Position:
    def __init__(self):
        self.left   = Inches(1)
        self.top    = Inches(1)
        self.width  = Inches(6)
        self.height = Inches(3)

def create_textbox(slide, text):
    '''
    Create new text box.
    [Text-related objects — python-pptx 0.6.21 documentation](https://python-pptx.readthedocs.io/en/latest/api/text.html#textframe-objects)
    '''
    pos        = Position()
    shapes     = slide.shapes
    txtbx      = shapes.add_textbox(pos.left, pos.top, pos.width, pos.height)
    txtbx.text = text


SLD_LAYOUT_TITLE_AND_CONTENT = 0

filename = 'test.pptx'
prs, file_exist = pptxBasics.open_pptx(filename)

# Set slide size
prs.slide_height=Inches(9)
prs.slide_width=Inches(16)

slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
slide = prs.slides.add_slide(slide_layout)

# slide = prs.slides[0]



# title = slide.shapes.title
subtitle = slide.placeholders[1]    
# title.text = "Hello, World2!"
subtitle.text = "python-pptx was here!"

position = [Cm(5.0), Cm(5.0), Cm(5.0), Cm(5.0)]
create_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, position)

# Add text box
text = "This is python-pptx text box object.\n Sample paragraph."
create_textbox(slide, text)

shapes = slide.shapes
create_date = "2021/09/26"
author      = "Fumihachi"
signature = shapes.add_textbox(Inches(2.5), Inches(7), Inches(1), Inches(1))
signature.text = create_date + "\n"+ author

# shapes.paragraphs[0].alignment = PP_ALIGN.LEFT
title = shapes.title
# title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
title.text = "Hello, World!"
# title.left = Inches(2.5)
# title.top  = Inches(2)

# Set title allignment
for sld in prs.slides:
    if sld.shapes.title:
        sld.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Save
pptxBasics.save_pptx(filename, prs, 0)
# pptxBasics.save_pptx(filename, prs, file_exist)





